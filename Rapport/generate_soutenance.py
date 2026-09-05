# -*- coding: utf-8 -*-
"""Génération du support PowerPoint de soutenance (15 min) — Alumni CRM.
Approche « documentation comme du code » : le deck est décrit en Python
et régénérable à tout moment. Design moderne : fond clair, bande latérale
indigo, cartes arrondies, encart « Exemple » sur chaque diapositive.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(OUTPUT_DIR, "..", "image")

# Palette moderne (indigo + vert émeraude, fond très clair)
BLUE = RGBColor(67, 56, 202)         # indigo vif
BLUE_DARK = RGBColor(37, 30, 122)     # indigo foncé
DARK = RGBColor(15, 23, 42)           # presque noir (slate-900)
GRAY = RGBColor(100, 116, 139)        # gris ardoise
LIGHT = RGBColor(238, 242, 255)       # indigo très clair
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(16, 185, 129)       # émeraude

SW, SH = Inches(13.333), Inches(7.5)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _side_bar(slide, color=BLUE):
    """Bande latérale indigo moderne sur le bord gauche."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.09), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def _accent_rule(slide, x, y, w=Inches(0.7), color=ACCENT):
    """Petite règle de couleur sous le titre."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.07))
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False


def _top_bar(slide, title, subtitle=None):
    """Titre moderne : texte + règle émeraude, sans gros bandeau plein."""
    _side_bar(slide)
    _text(slide, Inches(0.65), Inches(0.42), Inches(12.1), Inches(0.7),
          title, size=30, bold=True, color=BLUE_DARK)
    if subtitle:
        _text(slide, Inches(0.68), Inches(1.02), Inches(12.1), Inches(0.4),
              subtitle, size=14, color=GRAY)
    _accent_rule(slide, Inches(0.7), Inches(1.5))


def _box(slide, x, y, w, h, color=WHITE, line=None, round_=0.08):
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        bx.adjustments[0] = round_
    except Exception:
        pass
    bx.fill.solid()
    bx.fill.fore_color.rgb = color
    if line:
        bx.line.color.rgb = line
        bx.line.width = Pt(1)
    else:
        bx.line.fill.background()
    bx.shadow.inherit = False
    return bx


def _text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(2)
    return tb


def _bullets(slide, x, y, w, h, items, size=15, gap=7, bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.03)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ""
        p.font.name = "Calibri"
        r = p.add_run()
        r.text = "▸  "
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = bullet_color
        r.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(size)
        r2.font.color.rgb = DARK
        r2.font.name = "Calibri"
        p.space_after = Pt(gap)
    return tb


def _add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        return False
    pic = slide.shapes.add_picture(path, x, y, width=w)
    if h:
        pic.height = h
    return True


def _footer(slide, num):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Alumni CRM — Stage PreMSc 2026 — IONIS-STM        %d" % num
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.font.name = "Calibri"


def _screenshot(slide, path, left=Inches(1.9), top=Inches(1.35), ratio=1.6):
    h = Inches(5.55)
    w = int(h * ratio)
    _add_image(slide, path, left, top, w, h)


def _example(slide, label, text, y=Inches(6.12)):
    """Encart moderne « Exemple / Explication » en bas de diapositive."""
    bx = _box(slide, Inches(0.7), y, Inches(11.95), Inches(0.82), color=LIGHT, line=ACCENT, round_=0.12)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx.left, bx.top, Inches(0.09), bx.height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    stripe.shadow.inherit = False
    _text(slide, Inches(0.95), y + Inches(0.05), Inches(11.5), Inches(0.3),
          "EXEMPLE — " + label, size=13, bold=True, color=ACCENT)
    _text(slide, Inches(0.95), y + Inches(0.34), Inches(11.5), Inches(0.44),
          text, size=14, color=DARK)


def _numbered_rows(slide, y0, pairs, box_color, step=Inches(0.72), title=True):
    y = y0
    for code, txt in pairs:
        box = _box(slide, Inches(0.9), y, Inches(0.85), Inches(0.55), color=box_color, round_=0.25)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        _text(slide, Inches(2.0), y + Inches(0.02), Inches(10.4), Inches(0.5),
              txt, size=16, bold=(title and code.startswith("O")))
        y += step
    return y


def _card(slide, x, y, w, h, title, text, head_color=BLUE):
    """Carte moderne avec titre coloré et corps gris."""
    bg = _box(slide, x, y, w, h, color=WHITE, line=LIGHT, round_=0.1)
    _text(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.4),
          title, size=15, bold=True, color=head_color)
    _text(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), h - Inches(0.6),
          text, size=12.5, color=GRAY)
    return bg


def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ---------- 1. Titre ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _side_bar(s, BLUE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), Inches(13.333), Inches(0.14))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
    _text(s, Inches(1.2), Inches(0.85), Inches(10.9), Inches(0.5),
          "CONCEPTION ET DÉVELOPPEMENT D'UN ALUMNI CRM", size=33, bold=True, color=BLUE,
          align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(0.5),
          "Système de suivi du parcours étudiant et de valorisation du réseau des anciens",
          size=18, color=GRAY, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(2.6), Inches(10.9), Inches(0.9),
          "Rafik DJEMADI", size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(3.45), Inches(10.9), Inches(0.6),
          "Stage de substitution — Programme Pré-MSc 2026", size=17, color=GRAY,
          align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(4.05), Inches(10.9), Inches(0.5),
          "IONIS-STM — Ionis Education Group", size=17, color=BLUE, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(6.25), Inches(10.9), Inches(0.5),
          "Soutenance — 18 septembre 2026", size=15, color=GRAY, align=PP_ALIGN.CENTER)

    # ---------- 2. Plan ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Plan de la présentation")
    outline = [
        ("1", "Contexte et objectifs"),
        ("2", "Méthodologie et architecture"),
        ("3", "Modèle de données (MCD)"),
        ("4", "Conformité RGPD"),
        ("5", "Démonstration — Admin et Alumni"),
        ("6", "Indicateurs d'insertion"),
        ("7", "Bilan et perspectives"),
    ]
    y = Inches(1.45)
    for num, txt in outline:
        box = _box(s, Inches(0.9), y, Inches(0.85), Inches(0.62), color=BLUE, round_=0.2)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        _accent_rule(s, Inches(1.95), y + Inches(0.14), color=ACCENT)
        _text(s, Inches(2.15), y + Inches(0.07), Inches(10.3), Inches(0.5), txt, size=19, color=DARK)
        y += Inches(0.78)
    _footer(s, 2)

    # ---------- 3. Contexte & enjeux ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Contexte et enjeux")
    _bullets(s, Inches(0.9), Inches(1.3), Inches(11.6), Inches(2.3), [
        "IONIS-STM : formations Pré-MSc / MSc1 / MSc2 (développement, cyber, data, management, marketing digital).",
        "Plusieurs centaines de diplômés par an : suivre leur insertion est stratégique (pilotage, tutelle CTI/HCERES, animation du réseau).",
        "Sans outil centralisé : données dispersées, indicateurs à la main, réseau inactif, RGPD non formalisé.",
    ], size=14, gap=6)
    box = _box(s, Inches(0.9), Inches(4.3), Inches(11.6), Inches(1.05), color=BLUE, round_=0.1)
    _text(s, Inches(1.2), Inches(4.45), Inches(11.0), Inches(0.8),
          "Problématique : comment structurer la donnée alumni, produire des indicateurs "
          "d'insertion fiables et animer le réseau dans le respect du RGPD ?",
          size=16, bold=True, color=WHITE)
    _example(s, "Contexte",
             "retrouver les anciens d'une promotion repose aujourd'hui sur des fichiers épars "
             "et des relances manuelles — la base centrale remplace ce bricolage.")
    _footer(s, 3)

    # ---------- 4. Objectifs ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Objectifs du stage")
    goals = [
        ("O1", "Un CRM complet : espace admin + espace alumni"),
        ("O2", "Une base relationnelle SQL modélisée (MCD/MLD)"),
        ("O3", "Un tableau de bord avec indicateurs d'insertion"),
        ("O4", "Une conformité RGPD réelle (consentement, export, suppression, audit)"),
        ("O5", "L'import/export automatisé (Excel / CSV)"),
    ]
    _numbered_rows(s, Inches(1.5), goals, ACCENT, step=Inches(0.78))
    _example(s, "Objectif O3",
             "un directeur voit en une page le taux d'emploi, le salaire moyen et la répartition "
             "par secteur de la promotion 2026.")
    _footer(s, 4)

    # ---------- 5. Méthodologie ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Méthodologie et déroulement")
    cards = [
        ("Analyse", "Cahier des charges, étude des solutions, choix de la stack.", BLUE),
        ("Conception", "MCD/MLD (Looping), règles d'intégrité, schéma d'API.", ACCENT),
        ("Développement", "Backend FastAPI, frontend React, conformité RGPD.", BLUE_DARK),
        ("Consolidation", "Revue, documentation « comme du code », guide des processus.", ACCENT),
    ]
    pos = [(Inches(0.9), Inches(1.55)), (Inches(6.9), Inches(1.55)),
           (Inches(0.9), Inches(3.7)), (Inches(6.9), Inches(3.7))]
    for (c_title, c_text, col), (cx, cy) in zip(cards, pos):
        _card(s, cx, cy, Inches(5.6), Inches(2.0), c_title, c_text, head_color=col)
    _example(s, "Démarche itérative",
             "chaque brique est livrée puis testée sur les parcours réels (connexion OTP, "
             "import Excel, demande RGPD) avant d'enchaîner la suivante.")
    _footer(s, 5)

    # ---------- 6. Architecture 3-tiers ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Architecture 3-tiers")
    _bullets(s, Inches(0.9), Inches(1.35), Inches(6.0), Inches(3.0), [
        "Frontend React + Vite (SPA) : espace Admin et Alumni, proxy /api → 8000.",
        "Backend FastAPI : 83 endpoints, 16 routeurs, validations Pydantic, Swagger.",
        "PostgreSQL : 14 tables, 16 migrations versionnées, JSONB.",
        "Sécurité : OTP email, JWT, clé API admin, anti-IDOR.",
    ], size=14, gap=7)
    _example(s, "Flux d'une requête",
             "un alumni saisit son code OTP à 6 chiffres → FastAPI le valide → un JWT est "
             "émis pour accéder à son profil.", y=Inches(5.35))
    layers = [
        ("FRONTEND", "React + Vite\nSPA interactive\nproxy /api", BLUE),
        ("BACKEND", "FastAPI\n83 endpoints\nvalidations", ACCENT),
        ("BASE", "PostgreSQL\n14 tables\nJSONB", DARK),
    ]
    x = Inches(7.2)
    for nom, desc, col in layers:
        box = _box(s, x, Inches(1.55), Inches(1.85), Inches(3.3), color=col, round_=0.1)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = nom
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        for ln in desc.split("\n"):
            p2 = tf.add_paragraph()
            p2.text = ln
            p2.font.size = Pt(12)
            p2.font.color.rgb = WHITE
            p2.font.name = "Calibri"
            p2.alignment = PP_ALIGN.CENTER
        if x < Inches(11.0):
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.85), Inches(2.8), Inches(0.22), Inches(0.45))
            arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background(); arr.shadow.inherit = False
        x += Inches(2.12)
    _footer(s, 6)

    # ---------- 7. Modèle de données + capture MCD ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Modèle de données (MCD)")
    _bullets(s, Inches(0.9), Inches(1.25), Inches(11.6), Inches(1.0), [
        "5 domaines, 14 tables : étudiants, parcours pro, RGPD, questionnaires, INFRA.",
        "Relations clés : ETUDIANT → PROMOTION (N:1) ; EXPERIENCE_PRO (postes & salaires) ; OBTIENT (certifications N:M).",
    ], size=14, gap=5)
    mcd = os.path.join(FIG_DIR, "MCD.png")
    if os.path.exists(mcd):
        h_mcd = Inches(3.55)
        w_mcd = int(h_mcd * 1.972)
        left = int((SW - w_mcd) / 2)
        _add_image(s, mcd, left, Inches(2.25), w_mcd, h_mcd)
    _example(s, "Lecture du modèle",
             "un étudiant appartient à UNE promotion (N:1) mais peut obtenir plusieurs "
             "certifications (relation N:M OBTIENT) — la cardinalité guide les clés étrangères.")
    _footer(s, 7)

    # ---------- 8. RGPD ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Conformité RGPD")
    _bullets(s, Inches(0.9), Inches(1.3), Inches(11.6), Inches(1.5), [
        "Consentement explicite (4 types), horodaté et réellement consommé.",
        "Droits : export en auto-service (JSON/Excel/CSV), suppression avec workflow verrouillé et anonymisation.",
        "Traçabilité : journal d'audit, purge différée des comptes anonymisés.",
    ], size=14, gap=4)
    _add_image(s, os.path.join(FIG_DIR, "anC_consentement_light.png"), Inches(0.9), Inches(3.0), Inches(5.5))
    _add_image(s, os.path.join(FIG_DIR, "anB_demandes_rgpd_light.png"), Inches(6.7), Inches(3.0), Inches(5.5))
    _example(s, "Parcours de suppression",
             "l'alumni demande la suppression → l'admin verrouille la demande (anti-double traitement) "
             "→ décision traitee → compte anonymisé et tracé dans AUDIT_LOG.")
    _footer(s, 8)

    # ---------- 9. Démo Admin ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Démonstration — Espace Administration", "Tableau de bord, annuaire, RGPD")
    shot = os.path.join(FIG_DIR, "anB_dashboard_light.png")
    if os.path.exists(shot):
        h_demo = Inches(4.6)
        w_demo = int(h_demo * 1.6)
        _add_image(s, shot, Inches(0.6), Inches(1.7), w_demo, h_demo)
    _card(s, Inches(8.3), Inches(1.75), Inches(4.45), Inches(4.5), "Ce que vous verrez",
          "▸ KPI en haut : taux d'emploi, salaire moyen, anciens contactés.\n"
          "▸ Graphiques par secteur et par type de contrat.\n"
          "▸ Annuaire filtrable : promotion, secteur, consentement.\n"
          "▸ Demandes RGPD : prendre en charge, traiter, exporter.",
          head_color=BLUE)
    _footer(s, 9)

    # ---------- 10. Démo Alumni ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Démonstration — Espace Alumni", "Profil, parcours, consentement, questionnaire")
    shot = os.path.join(FIG_DIR, "anC_parcours_light.png")
    if os.path.exists(shot):
        h_demo = Inches(4.6)
        w_demo = int(h_demo * 1.6)
        _add_image(s, shot, Inches(0.6), Inches(1.7), w_demo, h_demo)
    _card(s, Inches(8.3), Inches(1.75), Inches(4.45), Inches(4.5), "Ce que vous verrez",
          "▸ Inscription avec OTP par e-mail (code à 6 chiffres).\n"
          "▸ Mon parcours : expériences, certifications, salaire.\n"
          "▸ Consentements RGPD en auto-service.\n"
          "▸ Questionnaire annuel pré-rempli.\n"
          "▸ Export de ses propres données.",
          head_color=ACCENT)
    _footer(s, 10)

    # ---------- 11. Indicateurs ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Indicateurs d'insertion professionnelle")
    _bullets(s, Inches(0.9), Inches(1.35), Inches(5.9), Inches(3.3), [
        "8 indicateurs reproductibles, formule et source explicites.",
        "Taux d'emploi à 6 mois (expériences actives).",
        "Adéquation formation/emploi via les tags KPI.",
        "Salaire moyen / min / max (champ annuel).",
        "Répartition par promotion et secteur.",
        "Pas de chiffre trompeur : cohortes immatures → « non disponible ».",
    ], size=14, gap=5)
    _example(s, "Calcul",
             "taux d'emploi à 6 mois = diplômés de la promo en poste actif à la date de "
             "référence / diplômés de la promo (cohorte mature).", y=Inches(5.65))
    _add_image(s, os.path.join(FIG_DIR, "anB_dashboard_light.png"), Inches(7.0), Inches(1.9), Inches(5.6))
    _footer(s, 11)

    # ---------- 12. Bilan & perspectives ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _top_bar(s, "Bilan et perspectives")
    _bullets(s, Inches(0.9), Inches(1.3), Inches(11.6), Inches(3.2), [
        "Prototype complet et opérationnel, documenté « comme du code » (rapports, cartographies, guide des processus, support de soutenance).",
        "Compétences : full-stack, sécurité applicative, migrations versionnées, RGPD.",
        "Perspectives : suite de tests automatisés, automatisation du questionnaire, frontend newsletter, mentorat, application mobile / PWA.",
    ], size=15, gap=8)
    _example(s, "Suite",
             "le dépôt Git versionne les migrations et les scripts de génération : "
             "le projet reste rejouable et transmissible à la reprise.")
    _footer(s, 12)

    # ---------- 13. Conclusion / fin ----------
    s = _blank(prs); _set_bg(s, WHITE)
    _side_bar(s, BLUE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
    _text(s, Inches(1.2), Inches(2.0), Inches(10.9), Inches(0.6),
          "Merci pour votre attention", size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(3.1), Inches(10.9), Inches(0.5),
          "Questions / échanges", size=22, color=GRAY, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.9),
          "Rafik DJEMADI\nAlumni CRM — Stage PreMSc 2026 — IONIS-STM", size=16, color=DARK,
          align=PP_ALIGN.CENTER)

    out = os.path.join(OUTPUT_DIR, "Soutenance - Alumni CRM - PreMSc 2026.pptx")
    prs.save(out)
    print("Soutenance PPTX genere (%d slides)." % len(prs.slides._sldIdLst))
    return out


if __name__ == "__main__":
    build()